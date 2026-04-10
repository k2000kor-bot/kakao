#!/usr/bin/env python3
"""
고급 파일 처리 및 지식 추출 엔진
- 다양한 파일 형식 지원 (PDF, DOCX, TXT, XLSX, PPT, 이미지, 오디오, 비디오)
- OCR 및 음성/영상 텍스트 추출
- 자동 지식 베이스 구축
- 글쓰기 소재 생성
"""

import os
import json
import sqlite3
import hashlib
import mimetypes
from datetime import datetime
from typing import List, Dict, Any, Optional, Union
from pathlib import Path
import asyncio
import logging

from fastapi import FastAPI, UploadFile, File, HTTPException, Form
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
import uvicorn

# 파일 처리 라이브러리
try:
    import PyPDF2
    import docx
    import openpyxl
    import pptx
    from PIL import Image
    import pytesseract
    import speech_recognition as sr
    import cv2
    import easyocr
except ImportError as e:
    print(f"선택적 패키지 누락: {e}. 기본 텍스트 처리만 지원됩니다.")

# 로깅 설정
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = FastAPI(
    title="고급 파일 처리 엔진",
    description="업로드된 파일에서 지식을 추출하여 글쓰기 소재로 활용하는 시스템",
    version="1.0.0"
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# 데이터 모델
class FileAnalysisResult(BaseModel):
    file_id: str
    file_name: str
    file_type: str
    file_size: int
    upload_time: datetime
    extracted_text: str
    key_topics: List[str]
    entities: Dict[str, List[str]]
    writing_materials: Dict[str, Any]
    knowledge_summary: str
    categorization: Dict[str, Any]
    metadata: Dict[str, Any]

class KnowledgeBase(BaseModel):
    project_id: str
    total_files: int
    total_knowledge_items: int
    categories: Dict[str, List[str]]
    key_concepts: List[str]
    writing_suggestions: List[str]
    cross_references: Dict[str, List[str]]

class WritingMaterial(BaseModel):
    id: str
    title: str
    content: str
    category: str
    keywords: List[str]
    source_files: List[str]
    confidence_score: float
    usage_suggestions: List[str]

class AdvancedFileProcessor:
    def __init__(self):
        self.storage_path = Path("uploads")
        self.storage_path.mkdir(exist_ok=True)
        self.init_database()
        
        # OCR 엔진 초기화
        try:
            self.ocr_reader = easyocr.Reader(['ko', 'en'])
            self.ocr_available = True
        except:
            self.ocr_available = False
            logger.warning("OCR 기능을 사용할 수 없습니다.")
        
        # 음성 인식 초기화
        try:
            self.speech_recognizer = sr.Recognizer()
            self.speech_available = True
        except:
            self.speech_available = False
            logger.warning("음성 인식 기능을 사용할 수 없습니다.")

    def init_database(self):
        """데이터베이스 초기화"""
        conn = sqlite3.connect('file_knowledge_base.db')
        cursor = conn.cursor()
        
        # 파일 분석 결과 테이블
        cursor.execute('''
            CREATE TABLE IF NOT EXISTS file_analysis (
                file_id TEXT PRIMARY KEY,
                project_id TEXT,
                file_name TEXT,
                file_type TEXT,
                file_size INTEGER,
                upload_time TIMESTAMP,
                extracted_text TEXT,
                key_topics TEXT,
                entities TEXT,
                writing_materials TEXT,
                knowledge_summary TEXT,
                categorization TEXT,
                metadata TEXT,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            )
        ''')
        
        # 지식 베이스 테이블
        cursor.execute('''
            CREATE TABLE IF NOT EXISTS knowledge_base (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                project_id TEXT,
                concept TEXT,
                description TEXT,
                category TEXT,
                keywords TEXT,
                source_files TEXT,
                confidence_score REAL,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            )
        ''')
        
        # 글쓰기 소재 테이블
        cursor.execute('''
            CREATE TABLE IF NOT EXISTS writing_materials (
                id TEXT PRIMARY KEY,
                project_id TEXT,
                title TEXT,
                content TEXT,
                category TEXT,
                keywords TEXT,
                source_files TEXT,
                confidence_score REAL,
                usage_suggestions TEXT,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            )
        ''')
        
        # 인덱스 생성
        cursor.execute('CREATE INDEX IF NOT EXISTS idx_project_id ON file_analysis(project_id)')
        cursor.execute('CREATE INDEX IF NOT EXISTS idx_file_type ON file_analysis(file_type)')
        cursor.execute('CREATE INDEX IF NOT EXISTS idx_kb_project ON knowledge_base(project_id)')
        cursor.execute('CREATE INDEX IF NOT EXISTS idx_wm_project ON writing_materials(project_id)')
        
        conn.commit()
        conn.close()

    async def process_file(self, file: UploadFile, project_id: str) -> FileAnalysisResult:
        """파일 처리 메인 함수"""
        logger.info(f"파일 처리 시작: {file.filename}")
        
        # 파일 ID 생성
        file_content = await file.read()
        file_id = hashlib.md5(f"{file.filename}_{len(file_content)}_{datetime.now()}".encode()).hexdigest()
        
        # 파일 저장
        file_path = self.storage_path / f"{file_id}_{file.filename}"
        with open(file_path, "wb") as f:
            f.write(file_content)
        
        # 파일 타입 감지
        file_type = self._detect_file_type(file.filename, file_content)
        
        # 텍스트 추출
        extracted_text = await self._extract_text(file_path, file_type)
        
        # 지식 분석
        analysis_result = await self._analyze_content(extracted_text, file.filename, file_type)
        
        # 결과 객체 생성
        result = FileAnalysisResult(
            file_id=file_id,
            file_name=file.filename,
            file_type=file_type,
            file_size=len(file_content),
            upload_time=datetime.now(),
            extracted_text=extracted_text,
            key_topics=analysis_result['topics'],
            entities=analysis_result['entities'],
            writing_materials=analysis_result['writing_materials'],
            knowledge_summary=analysis_result['summary'],
            categorization=analysis_result['categorization'],
            metadata=analysis_result['metadata']
        )
        
        # 데이터베이스 저장
        await self._save_analysis_result(result, project_id)
        
        # 지식 베이스 업데이트
        await self._update_knowledge_base(result, project_id)
        
        logger.info(f"파일 처리 완료: {file.filename}")
        return result

    def _detect_file_type(self, filename: str, content: bytes) -> str:
        """파일 타입 감지"""
        mime_type, _ = mimetypes.guess_type(filename)
        ext = Path(filename).suffix.lower()
        
        if ext in ['.pdf']:
            return 'pdf'
        elif ext in ['.docx', '.doc']:
            return 'document'
        elif ext in ['.xlsx', '.xls', '.csv']:
            return 'spreadsheet'
        elif ext in ['.pptx', '.ppt']:
            return 'presentation'
        elif ext in ['.txt', '.md']:
            return 'text'
        elif ext in ['.jpg', '.jpeg', '.png', '.bmp', '.tiff']:
            return 'image'
        elif ext in ['.mp3', '.wav', '.m4a', '.aac']:
            return 'audio'
        elif ext in ['.mp4', '.avi', '.mov', '.mkv']:
            return 'video'
        else:
            return 'unknown'

    async def _extract_text(self, file_path: Path, file_type: str) -> str:
        """파일에서 텍스트 추출"""
        try:
            if file_type == 'pdf':
                return await self._extract_pdf_text(file_path)
            elif file_type == 'document':
                return await self._extract_docx_text(file_path)
            elif file_type == 'spreadsheet':
                if file_path.suffix.lower() == '.csv':
                    return await self._extract_csv_text(file_path)
                return await self._extract_excel_text(file_path)
            elif file_type == 'presentation':
                return await self._extract_pptx_text(file_path)
            elif file_type == 'text':
                return await self._extract_plain_text(file_path)
            elif file_type == 'image':
                return await self._extract_image_text(file_path)
            elif file_type == 'audio':
                return await self._extract_audio_text(file_path)
            elif file_type == 'video':
                return await self._extract_video_text(file_path)
            else:
                return "텍스트를 추출할 수 없는 파일 형식입니다."
        except Exception as e:
            logger.error(f"텍스트 추출 오류: {e}")
            return f"텍스트 추출 중 오류 발생: {str(e)}"

    async def _extract_pdf_text(self, file_path: Path) -> str:
        """PDF에서 텍스트 추출"""
        try:
            import PyPDF2
            text = ""
            with open(file_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                for page in reader.pages:
                    text += page.extract_text() + "\n"
            return text.strip()
        except ImportError:
            return "PDF 처리 라이브러리가 설치되지 않았습니다."
        except Exception as e:
            return f"PDF 처리 오류: {str(e)}"

    async def _extract_docx_text(self, file_path: Path) -> str:
        """DOCX에서 텍스트 추출"""
        try:
            import docx
            doc = docx.Document(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text.strip()
        except ImportError:
            return "DOCX 처리 라이브러리가 설치되지 않았습니다."
        except Exception as e:
            return f"DOCX 처리 오류: {str(e)}"

    async def _extract_excel_text(self, file_path: Path) -> str:
        """Excel에서 텍스트 추출"""
        try:
            import openpyxl
            workbook = openpyxl.load_workbook(file_path)
            text = ""
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text += f"=== {sheet_name} ===\n"
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                    if row_text.strip():
                        text += row_text + "\n"
                text += "\n"
            return text.strip()
        except ImportError:
            return "Excel 처리 라이브러리가 설치되지 않았습니다."
        except Exception as e:
            return f"Excel 처리 오류: {str(e)}"

    async def _extract_csv_text(self, file_path: Path) -> str:
        """CSV에서 텍스트 추출 (엑셀보내기·한글 환경 cp949 대응)"""
        try:
            import csv
            from io import StringIO

            raw = file_path.read_bytes()
            text: Optional[str] = None
            for enc in ("utf-8-sig", "utf-8", "cp949", "euc-kr"):
                try:
                    text = raw.decode(enc)
                    break
                except UnicodeDecodeError:
                    continue
            if text is None:
                text = raw.decode("utf-8", errors="replace")

            out_lines: list[str] = []
            reader = csv.reader(StringIO(text))
            for row in reader:
                cells = ["" if c is None else str(c).strip() for c in row]
                line = " | ".join(cells)
                if line.strip():
                    out_lines.append(line)
            return "\n".join(out_lines).strip()
        except Exception as e:
            return f"CSV 처리 오류: {str(e)}"

    async def _extract_pptx_text(self, file_path: Path) -> str:
        """PowerPoint에서 텍스트 추출"""
        try:
            import pptx
            presentation = pptx.Presentation(file_path)
            text = ""
            for i, slide in enumerate(presentation.slides):
                text += f"=== 슬라이드 {i+1} ===\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                text += "\n"
            return text.strip()
        except ImportError:
            return "PowerPoint 처리 라이브러리가 설치되지 않았습니다."
        except Exception as e:
            return f"PowerPoint 처리 오류: {str(e)}"

    async def _extract_plain_text(self, file_path: Path) -> str:
        """일반 텍스트 파일에서 내용 추출"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except UnicodeDecodeError:
            try:
                with open(file_path, 'r', encoding='cp949') as file:
                    return file.read()
            except:
                with open(file_path, 'r', encoding='latin-1') as file:
                    return file.read()

    async def _extract_image_text(self, file_path: Path) -> str:
        """이미지에서 OCR로 텍스트 추출"""
        if not self.ocr_available:
            return "OCR 기능을 사용할 수 없습니다."
        
        try:
            result = self.ocr_reader.readtext(str(file_path))
            text = " ".join([item[1] for item in result])
            return text
        except Exception as e:
            return f"OCR 처리 오류: {str(e)}"

    async def _extract_audio_text(self, file_path: Path) -> str:
        """오디오에서 음성 인식으로 텍스트 추출"""
        if not self.speech_available:
            return "음성 인식 기능을 사용할 수 없습니다."
        
        try:
            # 실제 구현에서는 더 정교한 음성 인식 처리 필요
            return "음성 파일에서 텍스트 추출 기능은 추후 구현 예정입니다."
        except Exception as e:
            return f"음성 인식 오류: {str(e)}"

    async def _extract_video_text(self, file_path: Path) -> str:
        """비디오에서 텍스트 추출 (자막, OCR 등)"""
        try:
            return "비디오 파일에서 텍스트 추출 기능은 추후 구현 예정입니다."
        except Exception as e:
            return f"비디오 처리 오류: {str(e)}"

    async def _analyze_content(self, text: str, filename: str, file_type: str) -> Dict[str, Any]:
        """내용 분석 및 지식 추출"""
        # 키워드 및 토픽 추출
        topics = self._extract_topics(text)
        
        # 엔티티 추출
        entities = self._extract_entities(text)
        
        # 글쓰기 소재 생성
        writing_materials = self._generate_writing_materials(text, topics, entities)
        
        # 요약 생성
        summary = self._generate_summary(text)
        
        # 분류
        categorization = self._categorize_content(text, file_type)
        
        # 메타데이터
        metadata = {
            'word_count': len(text.split()),
            'char_count': len(text),
            'readability_score': self._calculate_readability(text),
            'language': self._detect_language(text),
            'processing_time': datetime.now().isoformat()
        }
        
        return {
            'topics': topics,
            'entities': entities,
            'writing_materials': writing_materials,
            'summary': summary,
            'categorization': categorization,
            'metadata': metadata
        }

    def _extract_topics(self, text: str) -> List[str]:
        """토픽 추출"""
        # 간단한 키워드 기반 토픽 추출
        topics = []
        
        # 한국어 핵심 키워드 패턴
        korean_patterns = [
            r'(?:재개발|재건축|아파트)',
            r'(?:시공사|건설사|건설)',
            r'(?:분석|검토|평가)',
            r'(?:조합원|주민|거주자)',
            r'(?:홍보|마케팅|광고)',
            r'(?:법적|법률|규정)',
            r'(?:가격|비용|예산)',
            r'(?:일정|스케줄|계획)'
        ]
        
        for pattern in korean_patterns:
            if re.search(pattern, text):
                topics.append(pattern.replace('(?:', '').replace('|', '/').replace(')', ''))
        
        # 영어 키워드도 포함
        english_keywords = ['project', 'analysis', 'development', 'construction', 'management']
        for keyword in english_keywords:
            if keyword.lower() in text.lower():
                topics.append(keyword)
        
        return list(set(topics))

    def _extract_entities(self, text: str) -> Dict[str, List[str]]:
        """엔티티 추출"""
        entities = {
            'organizations': [],
            'locations': [],
            'people': [],
            'dates': [],
            'numbers': []
        }
        
        # 조직명 패턴
        org_patterns = [
            r'((?:삼성|현대|대우|롯데|GS)(?:물산|건설|그룹|전자)?)',
            r'(\w+(?:조합|협회|회사|법인))',
        ]
        
        for pattern in org_patterns:
            matches = re.findall(pattern, text)
            entities['organizations'].extend(matches)
        
        # 지명 패턴
        location_patterns = [
            r'(개포\w*)',
            r'(\w+(?:구|동|시|군))',
        ]
        
        for pattern in location_patterns:
            matches = re.findall(pattern, text)
            entities['locations'].extend(matches)
        
        # 인명 패턴
        people_patterns = [
            r'([가-힣]{2,4})(?:씨|님|박사|교수|대표|이사|부장|과장)'
        ]
        
        for pattern in people_patterns:
            matches = re.findall(pattern, text)
            entities['people'].extend(matches)
        
        # 날짜 패턴
        date_patterns = [
            r'(\d{4}년\s*\d{1,2}월\s*\d{1,2}일)',
            r'(\d{4}-\d{1,2}-\d{1,2})',
            r'(\d{1,2}/\d{1,2}/\d{4})'
        ]
        
        for pattern in date_patterns:
            matches = re.findall(pattern, text)
            entities['dates'].extend(matches)
        
        # 숫자 패턴 (금액, 면적 등)
        number_patterns = [
            r'(\d+(?:,\d{3})*(?:원|만원|억원))',
            r'(\d+(?:\.\d+)?(?:평|㎡|m²))',
            r'(\d+(?:\.\d+)?%)'
        ]
        
        for pattern in number_patterns:
            matches = re.findall(pattern, text)
            entities['numbers'].extend(matches)
        
        # 중복 제거
        for key in entities:
            entities[key] = list(set(entities[key]))
        
        return entities

    def _generate_writing_materials(self, text: str, topics: List[str], entities: Dict[str, List[str]]) -> Dict[str, Any]:
        """글쓰기 소재 생성"""
        materials = {
            'key_points': [],
            'quotes': [],
            'statistics': [],
            'case_studies': [],
            'arguments': []
        }
        
        sentences = text.split('.')
        
        # 핵심 포인트 추출
        for sentence in sentences:
            sentence = sentence.strip()
            if len(sentence) > 20 and any(topic in sentence for topic in topics):
                materials['key_points'].append(sentence)
        
        # 인용문 추출
        quote_patterns = [
            r'"([^"]+)"',
            r"'([^']+)'",
            r'라고\s*(?:말했다|밝혔다|강조했다|설명했다)'
        ]
        
        for pattern in quote_patterns:
            matches = re.findall(pattern, text)
            materials['quotes'].extend(matches)
        
        # 통계 데이터 추출
        for number in entities['numbers']:
            context_start = max(0, text.find(number) - 50)
            context_end = min(len(text), text.find(number) + len(number) + 50)
            context = text[context_start:context_end]
            materials['statistics'].append({
                'value': number,
                'context': context.strip()
            })
        
        return materials

    def _generate_summary(self, text: str, max_length: int = 200) -> str:
        """텍스트 요약 생성"""
        sentences = [s.strip() for s in text.split('.') if len(s.strip()) > 10]
        
        if len(sentences) <= 3:
            return text
        
        # 간단한 추출식 요약 (첫 문장 + 중간 중요 문장 + 마지막 문장)
        summary_sentences = []
        
        if sentences:
            summary_sentences.append(sentences[0])  # 첫 문장
        
        if len(sentences) > 2:
            # 중간에서 가장 긴 문장 (정보가 많을 가능성)
            middle_sentences = sentences[1:-1]
            if middle_sentences:
                longest_sentence = max(middle_sentences, key=len)
                summary_sentences.append(longest_sentence)
        
        if len(sentences) > 1:
            summary_sentences.append(sentences[-1])  # 마지막 문장
        
        summary = '. '.join(summary_sentences)
        
        if len(summary) > max_length:
            summary = summary[:max_length] + "..."
        
        return summary

    def _categorize_content(self, text: str, file_type: str) -> Dict[str, Any]:
        """내용 분류"""
        categories = {
            'primary_category': 'general',
            'sub_categories': [],
            'confidence_scores': {}
        }
        
        # 카테고리 키워드 매핑
        category_keywords = {
            'construction': ['재개발', '재건축', '건설', '시공', '아파트'],
            'legal': ['법률', '규정', '계약', '조례', '법적'],
            'financial': ['비용', '예산', '가격', '투자', '수익'],
            'marketing': ['홍보', '마케팅', '광고', '브랜딩'],
            'management': ['관리', '운영', '계획', '일정', '조직'],
            'analysis': ['분석', '검토', '평가', '조사', '연구']
        }
        
        # 각 카테고리별 점수 계산
        for category, keywords in category_keywords.items():
            score = sum(1 for keyword in keywords if keyword in text)
            if score > 0:
                categories['sub_categories'].append(category)
                categories['confidence_scores'][category] = score / len(keywords)
        
        # 주 카테고리 결정
        if categories['confidence_scores']:
            primary = max(categories['confidence_scores'], key=categories['confidence_scores'].get)
            categories['primary_category'] = primary
        
        return categories

    def _calculate_readability(self, text: str) -> float:
        """가독성 점수 계산 (간단한 버전)"""
        if not text:
            return 0.0
        
        sentences = len([s for s in text.split('.') if s.strip()])
        words = len(text.split())
        chars = len(text)
        
        if sentences == 0:
            return 0.0
        
        avg_sentence_length = words / sentences
        avg_word_length = chars / words if words > 0 else 0
        
        # 간단한 가독성 공식 (낮을수록 읽기 쉬움)
        readability = (avg_sentence_length * 0.4) + (avg_word_length * 0.6)
        
        # 0-1 스케일로 정규화 (1에 가까울수록 읽기 쉬움)
        return max(0.0, min(1.0, 1.0 - (readability / 100)))

    def _detect_language(self, text: str) -> str:
        """언어 감지"""
        korean_chars = len(re.findall(r'[가-힣]', text))
        english_chars = len(re.findall(r'[a-zA-Z]', text))
        total_chars = korean_chars + english_chars
        
        if total_chars == 0:
            return 'unknown'
        
        korean_ratio = korean_chars / total_chars
        
        if korean_ratio > 0.7:
            return 'korean'
        elif korean_ratio < 0.3:
            return 'english'
        else:
            return 'mixed'

    async def _save_analysis_result(self, result: FileAnalysisResult, project_id: str):
        """분석 결과 저장"""
        conn = sqlite3.connect('file_knowledge_base.db')
        cursor = conn.cursor()
        
        cursor.execute('''
            INSERT OR REPLACE INTO file_analysis 
            (file_id, project_id, file_name, file_type, file_size, upload_time,
             extracted_text, key_topics, entities, writing_materials, 
             knowledge_summary, categorization, metadata)
            VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
        ''', (
            result.file_id, project_id, result.file_name, result.file_type,
            result.file_size, result.upload_time, result.extracted_text,
            json.dumps(result.key_topics, ensure_ascii=False),
            json.dumps(result.entities, ensure_ascii=False),
            json.dumps(result.writing_materials, ensure_ascii=False),
            result.knowledge_summary,
            json.dumps(result.categorization, ensure_ascii=False),
            json.dumps(result.metadata, ensure_ascii=False)
        ))
        
        conn.commit()
        conn.close()

    async def _update_knowledge_base(self, result: FileAnalysisResult, project_id: str):
        """지식 베이스 업데이트"""
        conn = sqlite3.connect('file_knowledge_base.db')
        cursor = conn.cursor()
        
        # 토픽을 지식 베이스에 추가
        for topic in result.key_topics:
            cursor.execute('''
                INSERT INTO knowledge_base 
                (project_id, concept, description, category, keywords, source_files, confidence_score)
                VALUES (?, ?, ?, ?, ?, ?, ?)
            ''', (
                project_id, topic, f"{topic} 관련 내용", 
                result.categorization.get('primary_category', 'general'),
                json.dumps(result.key_topics, ensure_ascii=False),
                json.dumps([result.file_id], ensure_ascii=False),
                0.8
            ))
        
        # 글쓰기 소재 저장
        for i, key_point in enumerate(result.writing_materials.get('key_points', [])[:5]):  # 최대 5개
            material_id = f"{result.file_id}_material_{i}"
            cursor.execute('''
                INSERT OR REPLACE INTO writing_materials 
                (id, project_id, title, content, category, keywords, source_files, 
                 confidence_score, usage_suggestions)
                VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)
            ''', (
                material_id, project_id,
                f"{result.file_name}에서 추출한 핵심 내용 {i+1}",
                key_point,
                result.categorization.get('primary_category', 'general'),
                json.dumps(result.key_topics, ensure_ascii=False),
                json.dumps([result.file_id], ensure_ascii=False),
                0.7,
                json.dumps(['서론 작성시 활용', '근거 자료로 인용', '예시로 활용'], ensure_ascii=False)
            ))
        
        conn.commit()
        conn.close()

    async def get_project_knowledge_base(self, project_id: str) -> KnowledgeBase:
        """프로젝트 지식 베이스 조회"""
        conn = sqlite3.connect('file_knowledge_base.db')
        cursor = conn.cursor()
        
        # 파일 통계
        cursor.execute('SELECT COUNT(*) FROM file_analysis WHERE project_id = ?', (project_id,))
        total_files = cursor.fetchone()[0]
        
        # 지식 항목 통계
        cursor.execute('SELECT COUNT(*) FROM knowledge_base WHERE project_id = ?', (project_id,))
        total_knowledge_items = cursor.fetchone()[0]
        
        # 카테고리별 분류
        cursor.execute('''
            SELECT category, GROUP_CONCAT(concept) 
            FROM knowledge_base WHERE project_id = ? 
            GROUP BY category
        ''', (project_id,))
        categories = {row[0]: row[1].split(',') if row[1] else [] for row in cursor.fetchall()}
        
        # 핵심 개념
        cursor.execute('''
            SELECT concept, confidence_score 
            FROM knowledge_base WHERE project_id = ? 
            ORDER BY confidence_score DESC LIMIT 10
        ''', (project_id,))
        key_concepts = [row[0] for row in cursor.fetchall()]
        
        # 글쓰기 제안
        cursor.execute('''
            SELECT title FROM writing_materials WHERE project_id = ? 
            ORDER BY confidence_score DESC LIMIT 5
        ''', (project_id,))
        writing_suggestions = [row[0] for row in cursor.fetchall()]
        
        conn.close()
        
        return KnowledgeBase(
            project_id=project_id,
            total_files=total_files,
            total_knowledge_items=total_knowledge_items,
            categories=categories,
            key_concepts=key_concepts,
            writing_suggestions=writing_suggestions,
            cross_references={}  # 추후 구현
        )

    async def get_writing_materials(self, project_id: str, category: Optional[str] = None) -> List[WritingMaterial]:
        """글쓰기 소재 조회"""
        conn = sqlite3.connect('file_knowledge_base.db')
        cursor = conn.cursor()
        
        query = '''
            SELECT id, title, content, category, keywords, source_files, 
                   confidence_score, usage_suggestions
            FROM writing_materials WHERE project_id = ?
        '''
        params = [project_id]
        
        if category:
            query += ' AND category = ?'
            params.append(category)
        
        query += ' ORDER BY confidence_score DESC'
        
        cursor.execute(query, params)
        rows = cursor.fetchall()
        
        materials = []
        for row in rows:
            materials.append(WritingMaterial(
                id=row[0],
                title=row[1],
                content=row[2],
                category=row[3],
                keywords=json.loads(row[4]) if row[4] else [],
                source_files=json.loads(row[5]) if row[5] else [],
                confidence_score=row[6],
                usage_suggestions=json.loads(row[7]) if row[7] else []
            ))
        
        conn.close()
        return materials

# 전역 프로세서 인스턴스
file_processor = AdvancedFileProcessor()

# API 엔드포인트
@app.post("/api/v1/upload-file", response_model=FileAnalysisResult)
async def upload_and_analyze_file(
    file: UploadFile = File(...),
    project_id: str = Form(...)
):
    """파일 업로드 및 분석"""
    try:
        if not file.filename:
            raise HTTPException(status_code=400, detail="파일명이 없습니다.")
        
        result = await file_processor.process_file(file, project_id)
        return result
        
    except Exception as e:
        logger.error(f"파일 업로드 처리 오류: {e}")
        raise HTTPException(status_code=500, detail=f"파일 처리 중 오류 발생: {str(e)}")

@app.get("/api/v1/knowledge-base/{project_id}", response_model=KnowledgeBase)
async def get_knowledge_base(project_id: str):
    """프로젝트 지식 베이스 조회"""
    try:
        knowledge_base = await file_processor.get_project_knowledge_base(project_id)
        return knowledge_base
    except Exception as e:
        logger.error(f"지식 베이스 조회 오류: {e}")
        raise HTTPException(status_code=500, detail=f"지식 베이스 조회 중 오류 발생: {str(e)}")

@app.get("/api/v1/writing-materials/{project_id}")
async def get_writing_materials(project_id: str, category: Optional[str] = None):
    """글쓰기 소재 조회"""
    try:
        materials = await file_processor.get_writing_materials(project_id, category)
        return materials
    except Exception as e:
        logger.error(f"글쓰기 소재 조회 오류: {e}")
        raise HTTPException(status_code=500, detail=f"글쓰기 소재 조회 중 오류 발생: {str(e)}")

@app.get("/api/v1/health")
async def health_check():
    """헬스 체크"""
    return {
        "status": "healthy",
        "service": "advanced_file_processor",
        "version": "1.0.0",
        "features": {
            "ocr_available": file_processor.ocr_available,
            "speech_available": file_processor.speech_available
        }
    }

@app.get("/")
async def root():
    """루트 엔드포인트"""
    return {
        "service": "고급 파일 처리 엔진",
        "version": "1.0.0",
        "capabilities": [
            "다양한 파일 형식 지원",
            "자동 텍스트 추출",
            "지식 베이스 구축",
            "글쓰기 소재 생성",
            "내용 분석 및 분류"
        ],
        "supported_formats": [
            "PDF", "DOCX", "XLSX", "PPTX", "TXT", 
            "이미지 (OCR)", "오디오 (음성인식)", "비디오"
        ],
        "endpoints": {
            "upload": "/api/v1/upload-file",
            "knowledge_base": "/api/v1/knowledge-base/{project_id}",
            "writing_materials": "/api/v1/writing-materials/{project_id}",
            "health": "/api/v1/health"
        }
    }

if __name__ == "__main__":
    _p = int(
        os.environ.get(
            "ADVANCED_FILE_PROCESSOR_PORT", os.environ.get("PORT", "8006")
        )
    )
    uvicorn.run(app, host="0.0.0.0", port=_p)
