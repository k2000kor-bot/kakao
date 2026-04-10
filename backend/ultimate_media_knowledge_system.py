#!/usr/bin/env python3
"""
궁극의 미디어 지식 활용 시스템 v1.0
- ChatGPT-5 수준의 미디어 이해 및 지식 추출
- 설득력 있는 콘텐츠 생성
- 실시간 학습 및 적응
- 멀티모달 지능형 분석
"""

import os
import json
import logging
import sqlite3
import subprocess
import tempfile
from datetime import datetime
from typing import Dict, List, Any, Optional
from pathlib import Path
from dataclasses import dataclass, asdict
from enum import Enum
import numpy as np

# AI/ML 라이브러리
try:
    import cv2
    import pytesseract
    import easyocr
    from transformers import pipeline
    from sentence_transformers import SentenceTransformer
except ImportError as e:
    print(f"AI 라이브러리 누락: {e}")

# FastAPI
from fastapi import FastAPI, UploadFile, File, HTTPException, Form, Query, Body
from fastapi.responses import JSONResponse, PlainTextResponse
from fastapi.middleware.cors import CORSMiddleware

import uvicorn

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class MediaType(Enum):
    IMAGE = "image"
    VIDEO = "video"
    AUDIO = "audio"
    DOCUMENT = "document"
    PRESENTATION = "presentation"
    SPREADSHEET = "spreadsheet"


class KnowledgeLevel(Enum):
    BASIC = "basic"
    INTERMEDIATE = "intermediate"
    ADVANCED = "advanced"
    EXPERT = "expert"


@dataclass
class ExtractedKnowledge:
    """추출된 지식"""

    content: str
    confidence: float
    knowledge_type: str
    entities: List[str]
    relationships: Dict[str, List[str]]
    insights: List[str]
    source_location: str
    timestamp: datetime


@dataclass
class PersuasiveContent:
    """설득력 있는 콘텐츠"""

    title: str
    content: str
    arguments: List[str]
    evidence: List[str]
    emotional_appeal: Dict[str, float]
    logical_structure: Dict[str, Any]
    target_audience: str
    persuasion_techniques: List[str]


class UltimateMediaKnowledgeSystem:
    """궁극의 미디어 지식 활용 시스템"""

    def __init__(self):
        self.knowledge_base = {}
        self.learning_history = []
        self.persuasion_patterns = self._load_persuasion_patterns()
        self.ai_models = self._initialize_ai_models()
        self.media_processors = self._initialize_media_processors()
        # DB 초기화
        self.db_path = os.path.join(str(Path.cwd()), "backend", "media_knowledge.db")
        self._init_db()
        # DB에서 기존 지식 로드
        self._load_from_db()

    def _initialize_ai_models(self) -> Dict[str, Any]:
        """AI 모델 초기화"""
        models = {}

        try:
            # OCR 모델 (라이브러리 존재 시에만 초기화)
            if "easyocr" in globals():
                try:
                    models["ocr"] = easyocr.Reader(["ko", "en"])
                except Exception as ocr_err:
                    logger.warning(f"OCR 모델 초기화 실패: {ocr_err}")
            else:
                logger.info("OCR 라이브러리 미설치: OCR 기능 비활성화")

            # 텍스트/감정/NER/임베딩 모델은 환경에 따라 무거울 수 있으므로 실패해도 계속 진행
            try:
                models["text_analyzer"] = pipeline(
                    "text-classification", model="klue/bert-base"
                )
            except Exception as e_text:
                logger.warning(f"텍스트 분석 모델 초기화 실패: {e_text}")

            try:
                models["sentiment"] = pipeline(
                    "sentiment-analysis", model="klue/bert-base"
                )
            except Exception as e_sent:
                logger.warning(f"감정 분석 모델 초기화 실패: {e_sent}")

            try:
                models["ner"] = pipeline("ner", model="klue/bert-base")
            except Exception as e_ner:
                logger.warning(f"NER 모델 초기화 실패: {e_ner}")

            try:
                models["sentence_transformer"] = SentenceTransformer(
                    "jhgan/ko-sroberta-multitask"
                )
            except Exception as e_st:
                logger.warning(f"문장 임베딩 모델 초기화 실패: {e_st}")

            logger.info("✅ AI 모델 초기화 단계 완료")

        except Exception as e:
            logger.warning(f"AI 모델 초기화 중 알 수 없는 오류: {e}")

        return models

    def _initialize_media_processors(self) -> Dict[str, Any]:
        """미디어 처리기 초기화"""
        return {
            "image": self._process_image,
            "video": self._process_video,
            "audio": self._process_audio,
            "document": self._process_document,
            "presentation": self._process_presentation,
            "spreadsheet": self._process_spreadsheet,
        }

    def _init_db(self) -> None:
        """SQLite DB 초기화 (테이블 생성)"""
        try:
            os.makedirs(os.path.dirname(self.db_path), exist_ok=True)
            with sqlite3.connect(self.db_path) as conn:
                cur = conn.cursor()
                cur.execute(
                    """
                    CREATE TABLE IF NOT EXISTS knowledge_items (
                        id INTEGER PRIMARY KEY AUTOINCREMENT,
                        project_id TEXT,
                        content TEXT,
                        confidence REAL,
                        knowledge_type TEXT,
                        entities TEXT,
                        relationships TEXT,
                        insights TEXT,
                        source_location TEXT,
                        timestamp TEXT
                    )
                    """
                )
                cur.execute(
                    """
                    CREATE TABLE IF NOT EXISTS learning_history (
                        id INTEGER PRIMARY KEY AUTOINCREMENT,
                        project_id TEXT,
                        session_id TEXT,
                        media_type TEXT,
                        learning_outcome TEXT,
                        confidence_score REAL,
                        entities_count INTEGER,
                        insights_count INTEGER,
                        timestamp TEXT
                    )
                    """
                )
                # 인덱스 생성
                cur.execute(
                    """
                    CREATE INDEX IF NOT EXISTS idx_knowledge_items_project_id
                    ON knowledge_items(project_id)
                    """
                )
                cur.execute(
                    """
                    CREATE INDEX IF NOT EXISTS idx_learning_history_project_id
                    ON learning_history(project_id)
                    """
                )
                conn.commit()
            logger.info("✅ SQLite DB 초기화 완료")
        except Exception as e:
            logger.warning(f"SQLite DB 초기화 실패: {e}")

    def _save_knowledge_to_db(
        self, project_id: str, knowledge: "ExtractedKnowledge"
    ) -> None:
        """추출 지식을 DB에 저장"""
        try:
            with sqlite3.connect(self.db_path) as conn:
                cur = conn.cursor()
                cur.execute(
                    """
                    INSERT INTO knowledge_items (
                        project_id, content, confidence, knowledge_type, entities, relationships, insights, source_location, timestamp
                    ) VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)
                    """,
                    (
                        project_id,
                        knowledge.content,
                        float(knowledge.confidence),
                        knowledge.knowledge_type,
                        json.dumps(knowledge.entities, ensure_ascii=False),
                        json.dumps(knowledge.relationships, ensure_ascii=False),
                        json.dumps(knowledge.insights, ensure_ascii=False),
                        knowledge.source_location,
                        knowledge.timestamp.isoformat(),
                    ),
                )
                conn.commit()
        except Exception as e:
            logger.warning(f"지식 DB 저장 실패: {e}")

    def _save_learning_event_to_db(
        self, project_id: str, knowledge: "ExtractedKnowledge"
    ) -> None:
        """학습 이벤트를 DB에 저장"""
        try:
            session_id = f"session_{project_id}_{int(datetime.now().timestamp())}"
            learning_outcome = (
                "success" if knowledge.confidence > 0.3 else "needs_review"
            )
            with sqlite3.connect(self.db_path) as conn:
                cur = conn.cursor()
                cur.execute(
                    """
                    INSERT INTO learning_history (
                        project_id, session_id, media_type, learning_outcome, confidence_score, entities_count, insights_count, timestamp
                    ) VALUES (?, ?, ?, ?, ?, ?, ?, ?)
                    """,
                    (
                        project_id,
                        session_id,
                        knowledge.knowledge_type,
                        learning_outcome,
                        float(knowledge.confidence),
                        int(len(knowledge.entities)),
                        int(len(knowledge.insights)),
                        knowledge.timestamp.isoformat(),
                    ),
                )
                conn.commit()
        except Exception as e:
            logger.warning(f"학습 이벤트 DB 저장 실패: {e}")

    def _load_from_db(self) -> None:
        """DB에서 지식/학습 이력을 메모리로 로드"""
        try:
            with sqlite3.connect(self.db_path) as conn:
                cur = conn.cursor()
                # 지식 항목 로드
                cur.execute(
                    """
                    SELECT project_id, content, confidence, knowledge_type, entities, relationships, insights, source_location, timestamp
                    FROM knowledge_items
                    ORDER BY id ASC
                    """
                )
                rows = cur.fetchall()
                for row in rows:
                    (
                        project_id,
                        content,
                        confidence,
                        knowledge_type,
                        entities_json,
                        relationships_json,
                        insights_json,
                        source_location,
                        ts,
                    ) = row
                    try:
                        knowledge = ExtractedKnowledge(
                            content=content or "",
                            confidence=float(confidence or 0.0),
                            knowledge_type=knowledge_type or "document",
                            entities=json.loads(entities_json or "[]"),
                            relationships=json.loads(relationships_json or "{}"),
                            insights=json.loads(insights_json or "[]"),
                            source_location=source_location or "unknown",
                            timestamp=(
                                datetime.fromisoformat(ts) if ts else datetime.now()
                            ),
                        )
                        if project_id not in self.knowledge_base:
                            self.knowledge_base[project_id] = []
                        self.knowledge_base[project_id].append(knowledge)
                    except Exception as e:
                        logger.warning(f"DB 지식 로드 항목 무시: {e}")

                # 학습 이력 로드(요약만)
                cur.execute(
                    """
                    SELECT project_id, session_id, media_type, learning_outcome, confidence_score, entities_count, insights_count, timestamp
                    FROM learning_history
                    ORDER BY id ASC
                    """
                )
                events = cur.fetchall()
                for ev in events:
                    try:
                        (
                            project_id,
                            session_id,
                            media_type,
                            learning_outcome,
                            confidence_score,
                            entities_count,
                            insights_count,
                            ts,
                        ) = ev
                        self.learning_history.append(
                            {
                                "project_id": project_id,
                                "session_id": session_id,
                                "media_type": media_type,
                                "learning_outcome": learning_outcome,
                                "confidence_score": float(confidence_score or 0.0),
                                "entities_count": int(entities_count or 0),
                                "insights_count": int(insights_count or 0),
                                "timestamp": ts,
                            }
                        )
                    except Exception as e:
                        logger.warning(f"DB 학습 이력 로드 항목 무시: {e}")
        except Exception as e:
            logger.warning(f"DB 로드 실패: {e}")

    def _clear_project_knowledge(self, project_id: str) -> int:
        """프로젝트 지식/이력을 DB와 메모리에서 삭제하고 삭제된 지식 개수 반환"""
        deleted_count = 0
        try:
            with sqlite3.connect(self.db_path) as conn:
                cur = conn.cursor()
                # 삭제될 지식 개수 파악
                cur.execute(
                    "SELECT COUNT(*) FROM knowledge_items WHERE project_id = ?",
                    (project_id,),
                )
                row = cur.fetchone()
                deleted_count = int(row[0]) if row else 0
                # DB 삭제
                cur.execute(
                    "DELETE FROM knowledge_items WHERE project_id = ?", (project_id,)
                )
                cur.execute(
                    "DELETE FROM learning_history WHERE project_id = ?", (project_id,)
                )
                conn.commit()
            # 메모리 갱신
            if project_id in self.knowledge_base:
                del self.knowledge_base[project_id]
            # learning_history 메모리는 요약이므로 필터링 제거
            self.learning_history = [
                e for e in self.learning_history if e.get("project_id") != project_id
            ]
        except Exception as e:
            logger.warning(f"프로젝트 지식 삭제 실패: {e}")
        return deleted_count

    def _load_persuasion_patterns(self) -> Dict[str, Any]:
        """설득 패턴 로드"""
        return {
            "ethos": {
                "credibility": ["전문성", "경험", "권위"],
                "trustworthiness": ["정직성", "신뢰성", "투명성"],
            },
            "pathos": {
                "emotional_appeal": ["공감", "동정", "희망", "두려움"],
                "storytelling": ["서사", "비유", "은유"],
            },
            "logos": {
                "logical_argument": ["사실", "통계", "논리"],
                "evidence": ["연구", "데이터", "사례"],
            },
        }

    async def analyze_media_file(
        self, file: UploadFile, project_id: str
    ) -> Dict[str, Any]:
        """미디어 파일 종합 분석"""
        try:
            # 파일 저장
            file_path = await self._save_uploaded_file(file)

            # 파일 타입 감지
            media_type = self._detect_media_type(file.filename, file.content_type)

            # 미디어 처리
            processor = self.media_processors.get(media_type.value)
            if not processor:
                raise HTTPException(
                    status_code=400,
                    detail=f"지원하지 않는 파일 타입: {media_type.value}",
                )

            # 지식 추출
            extracted_knowledge = await processor(file_path, file.filename)

            # 지식 베이스 업데이트
            await self._update_knowledge_base(project_id, extracted_knowledge)

            # 설득력 있는 콘텐츠 생성
            persuasive_content = await self._generate_persuasive_content(
                extracted_knowledge
            )

            return {
                "file_analysis": {
                    "file_name": file.filename,
                    "file_size": file.size,
                    "media_type": media_type.value,
                    "analysis_timestamp": datetime.now().isoformat(),
                },
                "extracted_knowledge": asdict(extracted_knowledge),
                "persuasive_content": asdict(persuasive_content),
                "knowledge_summary": await self._generate_knowledge_summary(
                    extracted_knowledge
                ),
                "learning_insights": await self._generate_learning_insights(
                    extracted_knowledge
                ),
            }

        except Exception as e:
            logger.error(f"미디어 파일 분석 실패: {e}")
            raise HTTPException(
                status_code=500, detail=f"파일 분석 중 오류 발생: {str(e)}"
            )

    async def _process_image(self, file_path: str, filename: str) -> ExtractedKnowledge:
        """이미지 처리 및 지식 추출"""
        try:
            # 이미지 로드
            image = cv2.imread(file_path)

            # OCR 텍스트 추출
            ocr_text = self._extract_text_from_image(image)

            # 이미지 분석
            image_analysis = await self._analyze_image_content(image)

            # 지식 추출
            knowledge = await self._extract_knowledge_from_text(ocr_text, "image")

            # 이미지 특화 지식 추가
            knowledge.insights.extend(image_analysis.get("insights", []))
            knowledge.entities.extend(image_analysis.get("entities", []))

            return knowledge

        except Exception as e:
            logger.error(f"이미지 처리 실패: {e}")
            raise

    async def _process_video(self, file_path: str, filename: str) -> ExtractedKnowledge:
        """비디오 처리 및 지식 추출"""
        try:
            # 비디오 프레임 추출
            frames = self._extract_video_frames(file_path)

            # 오디오 추출
            audio_text = await self._extract_audio_from_video(file_path)

            # 프레임별 분석
            frame_insights = []
            for frame in frames[:10]:  # 처음 10프레임만 분석
                frame_analysis = await self._analyze_image_content(frame)
                frame_insights.extend(frame_analysis.get("insights", []))

            # 통합 지식 추출
            combined_text = f"{audio_text}\n" + "\n".join(frame_insights)
            knowledge = await self._extract_knowledge_from_text(combined_text, "video")

            return knowledge

        except Exception as e:
            logger.error(f"비디오 처리 실패: {e}")
            raise

    async def _process_audio(self, file_path: str, filename: str) -> ExtractedKnowledge:
        """오디오 처리 및 지식 추출"""
        try:
            # 음성 인식
            audio_text = await self._extract_audio_text(file_path)

            # 지식 추출
            knowledge = await self._extract_knowledge_from_text(audio_text, "audio")

            return knowledge

        except Exception as e:
            logger.error(f"오디오 처리 실패: {e}")
            raise

    async def _process_document(
        self, file_path: str, filename: str
    ) -> ExtractedKnowledge:
        """문서 처리 및 지식 추출"""
        try:
            # 문서 텍스트 추출
            document_text = await self._extract_document_text(file_path)

            # 지식 추출
            knowledge = await self._extract_knowledge_from_text(
                document_text, "document"
            )

            return knowledge

        except Exception as e:
            logger.error(f"문서 처리 실패: {e}")
            raise

    async def _extract_knowledge_from_text(
        self, text: str, source_type: str
    ) -> ExtractedKnowledge:
        """텍스트에서 지식 추출"""
        try:
            # 엔터티 추출
            entities = await self._extract_entities(text)

            # 관계 추출
            relationships = await self._extract_relationships(text, entities)

            # 인사이트 생성
            insights = await self._generate_insights(text, entities)

            # 신뢰도 계산
            confidence = self._calculate_confidence(text, entities)

            return ExtractedKnowledge(
                content=text,
                confidence=confidence,
                knowledge_type=source_type,
                entities=entities,
                relationships=relationships,
                insights=insights,
                source_location=source_type,
                timestamp=datetime.now(),
            )

        except Exception as e:
            logger.error(f"지식 추출 실패: {e}")
            raise

    async def _generate_persuasive_content(
        self, knowledge: ExtractedKnowledge
    ) -> PersuasiveContent:
        """설득력 있는 콘텐츠 생성"""
        try:
            # 핵심 주제 추출
            main_topics = self._extract_main_topics(knowledge.content)

            # 논리적 구조 생성
            logical_structure = await self._create_logical_structure(knowledge)

            # 감정적 호소 생성
            emotional_appeal = await self._create_emotional_appeal(knowledge)

            # 설득 기법 적용
            persuasion_techniques = await self._apply_persuasion_techniques(knowledge)

            # 콘텐츠 생성
            content = await self._generate_content_with_persuasion(
                knowledge, logical_structure, emotional_appeal, persuasion_techniques
            )

            return PersuasiveContent(
                title=f"{main_topics[0]}에 대한 설득력 있는 분석",
                content=content,
                arguments=logical_structure.get("arguments", []),
                evidence=logical_structure.get("evidence", []),
                emotional_appeal=emotional_appeal,
                logical_structure=logical_structure,
                target_audience="일반 대중",
                persuasion_techniques=persuasion_techniques,
            )

        except Exception as e:
            logger.error(f"설득력 있는 콘텐츠 생성 실패: {e}")
            raise

    async def _create_logical_structure(
        self, knowledge: ExtractedKnowledge
    ) -> Dict[str, Any]:
        """논리적 구조 생성"""
        structure = {
            "introduction": [],
            "arguments": [],
            "evidence": [],
            "counter_arguments": [],
            "conclusion": [],
        }

        # 주요 논점 추출
        key_points = self._extract_key_points(knowledge.content)

        # 증거 추출
        evidence = self._extract_evidence(knowledge.content)

        # 반론 생성
        counter_arguments = await self._generate_counter_arguments(key_points)

        structure["arguments"] = key_points
        structure["evidence"] = evidence
        structure["counter_arguments"] = counter_arguments

        return structure

    async def _create_emotional_appeal(
        self, knowledge: ExtractedKnowledge
    ) -> Dict[str, float]:
        """감정적 호소 생성"""
        emotional_scores = {
            "공감": 0.0,
            "동정": 0.0,
            "희망": 0.0,
            "두려움": 0.0,
            "분노": 0.0,
            "기쁨": 0.0,
        }

        # 텍스트 감정 분석
        sentiment_result = self.ai_models.get("sentiment")
        if sentiment_result:
            try:
                result = sentiment_result(knowledge.content[:512])  # 첫 512자만 분석
                if result and len(result) > 0:
                    score = result[0]["score"]
                    label = result[0]["label"]

                    if "positive" in label.lower():
                        emotional_scores["희망"] = score
                        emotional_scores["기쁨"] = score
                    elif "negative" in label.lower():
                        emotional_scores["두려움"] = score
                        emotional_scores["분노"] = score
            except Exception as e:
                logger.warning(f"감정 분석 실패: {e}")

        return emotional_scores

    async def _apply_persuasion_techniques(
        self, knowledge: ExtractedKnowledge
    ) -> List[str]:
        """설득 기법 적용"""
        techniques = []

        # 논리적 설득 (Logos)
        if len(knowledge.entities) > 3:
            techniques.append("사실 기반 논증")

        if any("통계" in insight for insight in knowledge.insights):
            techniques.append("통계적 증거")

        # 감정적 설득 (Pathos)
        if any("사람" in entity for entity in knowledge.entities):
            techniques.append("인간적 호소")

        if any("가족" in entity or "아이" in entity for entity in knowledge.entities):
            techniques.append("가족적 감정")

        # 윤리적 설득 (Ethos)
        if any("전문가" in entity or "연구" in entity for entity in knowledge.entities):
            techniques.append("전문성 강조")

        if any("경험" in insight for insight in knowledge.insights):
            techniques.append("경험 기반 신뢰")

        return techniques

    async def _generate_content_with_persuasion(
        self,
        knowledge: ExtractedKnowledge,
        logical_structure: Dict[str, Any],
        emotional_appeal: Dict[str, float],
        persuasion_techniques: List[str],
    ) -> str:
        """설득력을 가진 콘텐츠 생성"""

        content_parts = []

        # 도입부
        intro = f"분석된 내용을 바탕으로 {knowledge.knowledge_type}에서 중요한 정보를 발견했습니다. "
        intro += "이 정보는 우리의 이해를 크게 향상시킬 수 있는 가치 있는 지식입니다."
        content_parts.append(intro)

        # 주요 논점
        if logical_structure.get("arguments"):
            content_parts.append("\n📋 **주요 논점:**")
            for i, argument in enumerate(logical_structure["arguments"][:3], 1):
                content_parts.append(f"{i}. {argument}")

        # 증거
        if logical_structure.get("evidence"):
            content_parts.append("\n🔍 **지원 증거:**")
            for evidence in logical_structure["evidence"][:3]:
                content_parts.append(f"• {evidence}")

        # 인사이트
        if knowledge.insights:
            content_parts.append("\n💡 **핵심 인사이트:**")
            for insight in knowledge.insights[:3]:
                content_parts.append(f"• {insight}")

        # 설득 기법 적용
        if persuasion_techniques:
            content_parts.append("\n🎯 **설득 요소:**")
            for technique in persuasion_techniques:
                content_parts.append(f"• {technique}")

        # 결론
        conclusion = "\n✨ **결론:** "
        conclusion += "이 분석을 통해 우리는 더 깊은 이해와 통찰력을 얻을 수 있습니다. "
        conclusion += "이러한 지식은 실용적이고 가치 있는 응용을 가능하게 합니다."
        content_parts.append(conclusion)

        return "\n".join(content_parts)

    def _extract_text_from_image(self, image: np.ndarray) -> str:
        """이미지에서 텍스트 추출"""
        try:
            # OCR 수행
            ocr_reader = self.ai_models.get("ocr")
            if ocr_reader:
                results = ocr_reader.readtext(image)
                texts = [
                    result[1] for result in results if result[2] > 0.5
                ]  # 신뢰도 50% 이상
                return " ".join(texts)
            else:
                # pytesseract 사용
                return pytesseract.image_to_string(image, lang="kor+eng")
        except Exception as e:
            logger.error(f"이미지 텍스트 추출 실패: {e}")
            return ""

    async def _analyze_image_content(self, image: np.ndarray) -> Dict[str, Any]:
        """이미지 내용 분석"""
        analysis = {"insights": [], "entities": [], "objects": []}

        try:
            # 객체 감지
            objects = self._detect_objects(image)
            analysis["objects"] = objects

            # 객체 기반 인사이트 생성
            for obj in objects:
                analysis["insights"].append(f"이미지에서 {obj}가 감지되었습니다.")
                analysis["entities"].append(obj)

            # 이미지 품질 분석
            quality_score = self._assess_image_quality(image)
            if quality_score > 0.8:
                analysis["insights"].append("이미지 품질이 우수합니다.")
            elif quality_score < 0.5:
                analysis["insights"].append("이미지 품질 개선이 필요합니다.")

        except Exception as e:
            logger.error(f"이미지 분석 실패: {e}")

        return analysis

    def _detect_objects(self, image: np.ndarray) -> List[str]:
        """객체 감지"""
        objects = []
        try:
            # 간단한 객체 감지 (실제로는 YOLO 등 사용)
            gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
            edges = cv2.Canny(gray, 50, 150)
            contours, _ = cv2.findContours(
                edges, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE
            )

            # 큰 객체만 선택
            for contour in contours:
                area = cv2.contourArea(contour)
                if area > 1000:  # 면적이 1000 이상인 객체
                    objects.append("감지된 객체")
                    break

        except Exception as e:
            logger.error(f"객체 감지 실패: {e}")

        return objects

    def _assess_image_quality(self, image: np.ndarray) -> float:
        """이미지 품질 평가"""
        try:
            # 밝기 분석
            brightness = np.mean(image)

            # 대비 분석
            contrast = np.std(image)

            # 선명도 분석
            laplacian_var = cv2.Laplacian(
                cv2.cvtColor(image, cv2.COLOR_BGR2GRAY), cv2.CV_64F
            ).var()

            # 종합 품질 점수
            quality_score = min(
                1.0, (brightness / 255 + contrast / 255 + laplacian_var / 1000) / 3
            )

            return quality_score

        except Exception as e:
            logger.error(f"이미지 품질 평가 실패: {e}")
            return 0.5

    async def _extract_entities(self, text: str) -> List[str]:
        """엔터티 추출"""
        entities = []
        try:
            ner_model = self.ai_models.get("ner")
            if ner_model:
                results = ner_model(text[:512])  # 첫 512자만 분석
                entities = [
                    result["word"] for result in results if result["score"] > 0.5
                ]
            else:
                # 간단한 엔터티 추출
                words = text.split()
                entities = [
                    word
                    for word in words
                    if len(word) > 2 and any(char.isupper() for char in word)
                ]

        except Exception as e:
            logger.error(f"엔터티 추출 실패: {e}")

        return entities

    async def _extract_relationships(
        self, text: str, entities: List[str]
    ) -> Dict[str, List[str]]:
        """관계 추출"""
        relationships = {}
        try:
            for entity in entities:
                relationships[entity] = []
                # 간단한 관계 추출 (실제로는 더 정교한 NLP 사용)
                sentences = text.split(".")
                for sentence in sentences:
                    if entity in sentence:
                        other_entities = [
                            e for e in entities if e != entity and e in sentence
                        ]
                        relationships[entity].extend(other_entities)

        except Exception as e:
            logger.error(f"관계 추출 실패: {e}")

        return relationships

    async def _generate_insights(self, text: str, entities: List[str]) -> List[str]:
        """인사이트 생성"""
        insights = []
        try:
            # 주요 주제 추출
            if entities:
                insights.append(f"주요 주제: {', '.join(entities[:3])}")

            # 텍스트 길이 기반 인사이트
            if len(text) > 1000:
                insights.append("상세한 내용이 포함된 문서입니다.")
            elif len(text) < 100:
                insights.append("간단한 정보가 포함된 문서입니다.")

            # 특정 키워드 기반 인사이트
            if any(keyword in text.lower() for keyword in ["중요", "핵심", "주요"]):
                insights.append("중요한 정보가 포함되어 있습니다.")

            if any(keyword in text.lower() for keyword in ["문제", "이슈", "쟁점"]):
                insights.append("문제점이나 이슈가 언급되어 있습니다.")

        except Exception as e:
            logger.error(f"인사이트 생성 실패: {e}")

        return insights

    def _calculate_confidence(self, text: str, entities: List[str]) -> float:
        """신뢰도 계산"""
        confidence = 0.5  # 기본값

        try:
            # 텍스트 길이 기반
            if len(text) > 500:
                confidence += 0.2
            elif len(text) > 100:
                confidence += 0.1

            # 엔터티 수 기반
            if len(entities) > 5:
                confidence += 0.2
            elif len(entities) > 2:
                confidence += 0.1

            # 텍스트 품질 기반
            if any(char.isdigit() for char in text):
                confidence += 0.1  # 숫자 포함

            if any(char.isupper() for char in text):
                confidence += 0.1  # 대문자 포함

        except Exception as e:
            logger.error(f"신뢰도 계산 실패: {e}")

        return min(confidence, 1.0)

    def _extract_main_topics(self, text: str) -> List[str]:
        """주요 주제 추출"""
        topics = []
        try:
            # 간단한 키워드 추출
            words = text.split()
            word_freq = {}
            for word in words:
                if len(word) > 2:
                    word_freq[word] = word_freq.get(word, 0) + 1

            # 빈도순 정렬
            sorted_words = sorted(word_freq.items(), key=lambda x: x[1], reverse=True)
            topics = [word for word, freq in sorted_words[:5]]

        except Exception as e:
            logger.error(f"주제 추출 실패: {e}")

        return topics

    def _extract_key_points(self, text: str) -> List[str]:
        """핵심 포인트 추출"""
        points = []
        try:
            sentences = text.split(".")
            for sentence in sentences:
                sentence = sentence.strip()
                if len(sentence) > 20 and any(
                    keyword in sentence for keyword in ["중요", "핵심", "주요", "결론"]
                ):
                    points.append(sentence)

            # 최대 3개까지만 반환
            return points[:3]

        except Exception as e:
            logger.error(f"핵심 포인트 추출 실패: {e}")
            return []

    def _extract_evidence(self, text: str) -> List[str]:
        """증거 추출"""
        evidence = []
        try:
            sentences = text.split(".")
            for sentence in sentences:
                sentence = sentence.strip()
                if any(
                    keyword in sentence
                    for keyword in ["데이터", "통계", "연구", "조사", "결과"]
                ):
                    evidence.append(sentence)

            return evidence[:3]

        except Exception as e:
            logger.error(f"증거 추출 실패: {e}")
            return []

    async def _generate_counter_arguments(self, arguments: List[str]) -> List[str]:
        """반론 생성"""
        counter_arguments = []
        try:
            for argument in arguments:
                # 간단한 반론 생성
                if "중요" in argument:
                    counter_arguments.append(
                        "다른 관점에서 보면 덜 중요할 수 있습니다."
                    )
                elif "문제" in argument:
                    counter_arguments.append(
                        "이 문제에 대한 다른 해결책이 있을 수 있습니다."
                    )
                else:
                    counter_arguments.append(
                        "이 주장에 대한 반대 의견이 있을 수 있습니다."
                    )

            return counter_arguments[:3]

        except Exception as e:
            logger.error(f"반론 생성 실패: {e}")
            return []

    async def _generate_knowledge_summary(self, knowledge: ExtractedKnowledge) -> str:
        """지식 요약 생성"""
        summary = "📚 **지식 요약**\n\n"
        summary += f"**추출된 지식 타입:** {knowledge.knowledge_type}\n"
        summary += f"**신뢰도:** {knowledge.confidence:.2f}\n"
        summary += f"**발견된 엔터티:** {len(knowledge.entities)}개\n"
        summary += f"**핵심 인사이트:** {len(knowledge.insights)}개\n\n"

        if knowledge.entities:
            summary += f"**주요 엔터티:** {', '.join(knowledge.entities[:5])}\n"

        if knowledge.insights:
            summary += "**주요 인사이트:**\n"
            for insight in knowledge.insights[:3]:
                summary += f"• {insight}\n"

        return summary

    async def _generate_learning_insights(self, knowledge: ExtractedKnowledge) -> str:
        """학습 인사이트 생성"""
        insights = "🧠 **학습 인사이트**\n\n"

        # 지식 수준 평가
        if knowledge.confidence > 0.8:
            insights += "✅ **고품질 지식:** 신뢰도가 높은 우수한 정보입니다.\n"
        elif knowledge.confidence > 0.6:
            insights += "⚠️ **중간 품질 지식:** 추가 검증이 필요할 수 있습니다.\n"
        else:
            insights += "❌ **낮은 품질 지식:** 신중한 검토가 필요합니다.\n"

        # 활용 방안 제시
        insights += "\n**활용 방안:**\n"
        insights += "• 설득력 있는 논증에 활용\n"
        insights += "• 전문적 분석 자료로 활용\n"
        insights += "• 교육 자료로 활용\n"
        insights += "• 의사결정 지원 자료로 활용\n"

        return insights

    async def _save_uploaded_file(self, file: UploadFile) -> str:
        """업로드된 파일 저장"""
        upload_dir = Path("uploads")
        upload_dir.mkdir(exist_ok=True)

        file_path = upload_dir / file.filename
        with open(file_path, "wb") as buffer:
            content = await file.read()
            buffer.write(content)

        return str(file_path)

    def _detect_media_type(self, filename: str, content_type: str) -> MediaType:
        """미디어 타입 감지"""
        filename_lower = filename.lower()

        if any(
            ext in filename_lower for ext in [".jpg", ".jpeg", ".png", ".gif", ".bmp"]
        ):
            return MediaType.IMAGE
        elif any(ext in filename_lower for ext in [".mp4", ".avi", ".mov", ".wmv"]):
            return MediaType.VIDEO
        elif any(ext in filename_lower for ext in [".mp3", ".wav", ".aac", ".m4a"]):
            return MediaType.AUDIO
        elif any(ext in filename_lower for ext in [".pdf", ".doc", ".docx", ".txt"]):
            return MediaType.DOCUMENT
        elif any(ext in filename_lower for ext in [".ppt", ".pptx"]):
            return MediaType.PRESENTATION
        elif any(ext in filename_lower for ext in [".xls", ".xlsx", ".csv"]):
            return MediaType.SPREADSHEET
        else:
            return MediaType.DOCUMENT  # 기본값

    async def _update_knowledge_base(
        self, project_id: str, knowledge: ExtractedKnowledge
    ):
        """지식 베이스 업데이트"""
        if project_id not in self.knowledge_base:
            self.knowledge_base[project_id] = []

        self.knowledge_base[project_id].append(knowledge)

        # 학습 히스토리에 추가
        self.learning_history.append(
            {
                "project_id": project_id,
                "knowledge": asdict(knowledge),
                "timestamp": datetime.now().isoformat(),
            }
        )
        # DB에 영속화
        self._save_knowledge_to_db(project_id, knowledge)
        self._save_learning_event_to_db(project_id, knowledge)

    def _transcribe_audio_bytes(self, audio_bytes: bytes, mime: str = "audio/wav") -> Optional[str]:
        """오디오 바이트에서 텍스트 추출 (Whisper). 미설치 시 None."""
        try:
            import whisper
        except ImportError:
            logger.debug("Whisper 미설치, 오디오 텍스트 추출 스킵")
            return None
        suffix = ".wav" if "wav" in mime else ".mp3"
        with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as f:
            f.write(audio_bytes)
            path = f.name
        try:
            model = whisper.load_model("base")
            result = model.transcribe(path, language="ko", fp16=False)
            return (result.get("text") or "").strip() or None
        except Exception as e:
            logger.debug("Whisper transcribe 실패: %s", e)
            return None
        finally:
            try:
                os.unlink(path)
            except OSError:
                pass

    def _extract_video_frames(self, file_path: str) -> List[np.ndarray]:
        """비디오에서 프레임 추출"""
        frames = []
        try:
            cap = cv2.VideoCapture(file_path)
            frame_count = 0
            while cap.isOpened() and frame_count < 10:
                ret, frame = cap.read()
                if ret:
                    frames.append(frame)
                    frame_count += 1
                else:
                    break
            cap.release()
        except Exception as e:
            logger.error(f"비디오 프레임 추출 실패: {e}")
        return frames

    async def _extract_audio_from_video(self, file_path: str) -> str:
        """비디오에서 오디오 추출 후 텍스트 변환 (ffmpeg + Whisper)"""
        try:
            with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as tmp:
                out_path = tmp.name
            try:
                proc = subprocess.run(
                    [
                        "ffmpeg", "-y", "-i", file_path,
                        "-vn", "-acodec", "pcm_s16le", "-ar", "16000", "-ac", "1",
                        out_path,
                    ],
                    capture_output=True,
                    timeout=120,
                )
                if proc.returncode != 0 or not os.path.exists(out_path):
                    logger.warning("ffmpeg 오디오 추출 실패: %s", proc.stderr.decode(errors="ignore")[:200])
                    return "비디오에서 추출된 오디오 텍스트"
                with open(out_path, "rb") as f:
                    audio_bytes = f.read()
                text = self._transcribe_audio_bytes(audio_bytes, "audio/wav")
                return text if text else "비디오에서 추출된 오디오 텍스트"
            finally:
                try:
                    os.unlink(out_path)
                except OSError:
                    pass
        except FileNotFoundError:
            logger.debug("ffmpeg 미설치")
            return "비디오에서 추출된 오디오 텍스트"
        except subprocess.TimeoutExpired:
            logger.warning("ffmpeg 타임아웃")
            return "비디오에서 추출된 오디오 텍스트"
        except Exception as e:
            logger.warning("비디오 오디오 추출 실패: %s", e)
            return "비디오에서 추출된 오디오 텍스트"

    async def _extract_audio_text(self, file_path: str) -> str:
        """오디오에서 텍스트 추출 (Whisper)"""
        try:
            ext = Path(file_path).suffix.lower()
            mime_map = {
                ".wav": "audio/wav", ".mp3": "audio/mpeg", ".m4a": "audio/mp4",
                ".flac": "audio/flac", ".ogg": "audio/ogg", ".webm": "audio/webm",
            }
            mime = mime_map.get(ext, "audio/wav")
            with open(file_path, "rb") as f:
                audio_bytes = f.read()
            text = self._transcribe_audio_bytes(audio_bytes, mime)
            return text if text else "오디오에서 추출된 텍스트"
        except Exception as e:
            logger.warning("오디오 텍스트 추출 실패: %s", e)
            return "오디오에서 추출된 텍스트"

    async def _extract_document_text(self, file_path: str) -> str:
        """문서에서 텍스트 추출 (txt/pdf/docx 기본 지원)"""
        try:
            ext = Path(file_path).suffix.lower()
            # 단순 텍스트 계열
            if ext in {".txt", ".md", ".rtf"}:
                try:
                    return Path(file_path).read_text(encoding="utf-8")
                except Exception:
                    return Path(file_path).read_text(encoding="cp949", errors="ignore")

            # PDF
            if ext == ".pdf":
                try:
                    import PyPDF2  # type: ignore

                    text_parts: List[str] = []
                    with open(file_path, "rb") as f:
                        reader = PyPDF2.PdfReader(f)
                        for page in reader.pages:
                            extracted = page.extract_text() or ""
                            if extracted:
                                text_parts.append(extracted)
                    return "\n".join(text_parts)
                except Exception as e:
                    logger.warning(f"PDF 텍스트 추출 실패: {e}")

            # DOCX
            if ext == ".docx":
                try:
                    import docx  # type: ignore

                    document = docx.Document(file_path)
                    return "\n".join([p.text for p in document.paragraphs if p.text])
                except Exception as e:
                    logger.warning(f"DOCX 텍스트 추출 실패: {e}")

            # PPTX
            if ext in {".pptx", ".ppt"}:
                try:
                    from pptx import Presentation  # python-pptx

                    prs = Presentation(file_path)
                    parts = []
                    for i, slide in enumerate(prs.slides):
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text:
                                parts.append(shape.text)
                    return "\n".join(parts) if parts else ""
                except ImportError:
                    logger.debug("python-pptx 미설치")
                except Exception as e:
                    logger.warning(f"PPTX 텍스트 추출 실패: {e}")

            # CSV (텍스트로 직접 읽기)
            if ext == ".csv":
                try:
                    return Path(file_path).read_text(encoding="utf-8", errors="ignore")
                except Exception:
                    return Path(file_path).read_text(encoding="cp949", errors="ignore")

            # XLSX
            if ext == ".xlsx":
                try:
                    import openpyxl  # type: ignore

                    wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
                    parts = []
                    for sheet in wb.worksheets:
                        for row in sheet.iter_rows(values_only=True):
                            cells = [str(c) for c in row if c is not None and str(c).strip()]
                            if cells:
                                parts.append(" ".join(cells))
                    wb.close()
                    return "\n".join(parts) if parts else ""
                except ImportError:
                    logger.debug("openpyxl 미설치")
                except Exception as e:
                    logger.warning(f"엑셀 텍스트 추출 실패: {e}")

            # 기타는 텍스트 시도로 폴백
            try:
                return Path(file_path).read_text(encoding="utf-8", errors="ignore")
            except Exception:
                return Path(file_path).read_text(encoding="cp949", errors="ignore")

        except Exception as e:
            logger.error(f"문서 텍스트 추출 실패: {e}")
            return ""

    async def _process_presentation(
        self, file_path: str, filename: str
    ) -> ExtractedKnowledge:
        """프레젠테이션 처리 (pptx 슬라이드 텍스트 추출)"""
        text = await self._extract_document_text(file_path)
        return await self._extract_knowledge_from_text(text, "presentation")

    async def _process_spreadsheet(
        self, file_path: str, filename: str
    ) -> ExtractedKnowledge:
        """스프레드시트 처리 (xlsx/csv 텍스트 추출)"""
        text = await self._extract_document_text(file_path)
        return await self._extract_knowledge_from_text(text, "spreadsheet")


# FastAPI 앱 설정
app = FastAPI(
    title="궁극의 미디어 지식 활용 시스템",
    description="ChatGPT-5 수준의 미디어 이해 및 지식 추출 시스템",
    version="1.0.0",
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# 시스템 인스턴스
system = UltimateMediaKnowledgeSystem()


@app.post("/api/v1/analyze-media")
async def analyze_media_file(file: UploadFile = File(...), project_id: str = Form(...)):
    """미디어 파일 분석 및 지식 추출"""
    return await system.analyze_media_file(file, project_id)


@app.get("/api/v1/knowledge-base/{project_id}")
async def get_knowledge_base(project_id: str):
    """프로젝트 지식 베이스 조회"""
    try:
        # 메모리 우선, 없으면 DB 조회
        if project_id in system.knowledge_base:
            kb = system.knowledge_base[project_id]
            return {
                "project_id": project_id,
                "knowledge_items": len(kb),
                "knowledge_base": [asdict(k) for k in kb],
            }

        # DB에서 가져오기
        results = []
        with sqlite3.connect(system.db_path) as conn:
            cur = conn.cursor()
            cur.execute(
                """
                SELECT content, confidence, knowledge_type, entities, relationships, insights, source_location, timestamp
                FROM knowledge_items
                WHERE project_id = ?
                ORDER BY id DESC
                """,
                (project_id,),
            )
            rows = cur.fetchall()
            for row in rows:
                (
                    content,
                    confidence,
                    knowledge_type,
                    entities_json,
                    relationships_json,
                    insights_json,
                    source_location,
                    ts,
                ) = row
                results.append(
                    {
                        "content": content,
                        "confidence": float(confidence or 0.0),
                        "knowledge_type": knowledge_type,
                        "entities": json.loads(entities_json or "[]"),
                        "relationships": json.loads(relationships_json or "{}"),
                        "insights": json.loads(insights_json or "[]"),
                        "source_location": source_location,
                        "timestamp": ts,
                    }
                )
        return {
            "project_id": project_id,
            "knowledge_items": len(results),
            "knowledge_base": results,
        }
    except Exception as e:
        logger.error(f"지식 베이스 조회 실패: {e}")
        raise HTTPException(status_code=500, detail="지식 베이스 조회 실패")


@app.get("/api/v1/learning-history/{project_id}")
async def get_learning_history(project_id: str):
    """프로젝트별 학습 히스토리 조회"""
    try:
        sessions = []
        # DB에서 세션 조회
        with sqlite3.connect(system.db_path) as conn:
            cur = conn.cursor()
            cur.execute(
                """
                SELECT session_id, timestamp, project_id, media_type, learning_outcome, confidence_score, entities_count, insights_count
                FROM learning_history
                WHERE project_id = ?
                ORDER BY id DESC
                """,
                (project_id,),
            )
            rows = cur.fetchall()
            for row in rows:
                session_id, ts, pid, media_type, outcome, conf, ent, ins = row
                sessions.append(
                    {
                        "session_id": session_id,
                        "timestamp": ts,
                        "project_id": pid,
                        "media_type": media_type,
                        "learning_outcome": outcome,
                        "confidence_score": float(conf or 0.0),
                        "entities_count": int(ent or 0),
                        "insights_count": int(ins or 0),
                    }
                )

        return {
            "project_id": project_id,
            "total_events": len(sessions),
            "learning_sessions": sessions,
        }
    except Exception as e:
        logger.error(f"학습 히스토리 조회 실패: {e}")
        raise HTTPException(status_code=500, detail="학습 히스토리 조회 실패")


@app.delete("/api/v1/knowledge-base/{project_id}")
async def clear_project_knowledge(project_id: str):
    """프로젝트 지식/학습 이력 초기화"""
    try:
        deleted = system._clear_project_knowledge(project_id)
        return {"project_id": project_id, "deleted_items": deleted, "status": "cleared"}
    except Exception as e:
        logger.error(f"지식 초기화 실패: {e}")
        raise HTTPException(status_code=500, detail="지식 초기화 실패")


@app.get("/api/v1/knowledge-export/{project_id}")
async def export_knowledge(
    project_id: str,
    format: str = Query("json", description="json|csv"),
    q: str = Query(""),
    min_confidence: float = Query(0.0),
    media_type: Optional[str] = Query(None),
    start: Optional[str] = Query(None),
    end: Optional[str] = Query(None),
    order: str = Query("desc"),
    limit: int = Query(1000),
):
    """프로젝트 지식 베이스를 JSON 또는 CSV로 내보내기"""
    try:
        params = {
            "project_id": project_id,
            "q": q,
            "min_confidence": min_confidence,
            "media_type": media_type,
            "limit": limit,
            "start": start,
            "end": end,
            "order": order,
        }
        result = await search_knowledge(**params)  # type: ignore[arg-type]

        if format.lower() == "csv":
            headers = [
                "timestamp",
                "knowledge_type",
                "confidence",
                "entities_count",
                "insights_count",
            ]
            lines = [",".join(headers)]
            for k in result["matches"]:
                ts = k.get("timestamp", "")
                kt = k.get("knowledge_type", "")
                conf = str(k.get("confidence", 0))
                ec = str(len(k.get("entities", []) or []))
                ic = str(len(k.get("insights", []) or []))
                row = [
                    str(ts).replace("\n", " "),
                    str(kt).replace("\n", " "),
                    conf,
                    ec,
                    ic,
                ]
                lines.append(",".join([s.replace('"', '""') for s in row]))
            csv_text = "\n".join(lines)
            return PlainTextResponse(
                content=csv_text,
                media_type="text/csv; charset=utf-8",
                headers={
                    "Content-Disposition": f"attachment; filename=knowledge_{project_id}.csv"
                },
            )

        return JSONResponse(
            content=result,
            headers={
                "Content-Disposition": f"attachment; filename=knowledge_{project_id}.json"
            },
        )
    except Exception as e:
        logger.error(f"지식 내보내기 실패: {e}")
        raise HTTPException(status_code=500, detail="지식 내보내기 실패")


@app.get("/api/v1/health")
async def health_check():
    """시스템 상태 확인"""
    return {
        "status": "healthy",
        "version": "1.0.0",
        "ai_models_loaded": len(system.ai_models),
        "timestamp": datetime.now().isoformat(),
    }


@app.get("/api/v1/search-knowledge")
async def search_knowledge(
    project_id: str = Query(...),
    q: str = Query(""),
    min_confidence: float = Query(0.0),
    media_type: Optional[str] = Query(None),
    limit: int = Query(50),
    start: Optional[str] = Query(None, description="ISO8601 시작일시"),
    end: Optional[str] = Query(None, description="ISO8601 종료일시"),
    order: str = Query("desc", description="정렬: asc|desc"),
):
    """프로젝트 지식 베이스에서 키워드 검색"""
    try:
        results = []
        # 메모리에서 우선 검색
        in_mem = system.knowledge_base.get(project_id, [])
        for k in in_mem:
            text_blob = " ".join(
                [
                    k.content or "",
                    " ".join(k.entities or []),
                    " ".join(k.insights or []),
                ]
            )
            if q.lower() in text_blob.lower():
                ts = (
                    k.timestamp.isoformat()
                    if hasattr(k.timestamp, "isoformat")
                    else str(k.timestamp)
                )
                in_range = True
                if start and ts < start:
                    in_range = False
                if end and ts > end:
                    in_range = False
                if in_range and k.confidence >= float(min_confidence) and (
                    not media_type or k.knowledge_type == media_type
                ):
                    results.append(asdict(k))

        # DB에서도 검색 (보완)
        with sqlite3.connect(system.db_path) as conn:
            conn.create_function("lower", 1, lambda s: (s or "").lower())
            cur = conn.cursor()
            like = f"%{q.lower()}%"
            where = ["project_id = ?", "confidence >= ?", "(lower(content) LIKE ? OR lower(entities) LIKE ? OR lower(insights) LIKE ?)"]
            params: list = [project_id, float(min_confidence), like, like, like]
            if media_type:
                where.append("knowledge_type = ?")
                params.append(media_type)
            if start:
                where.append("timestamp >= ?")
                params.append(start)
            if end:
                where.append("timestamp <= ?")
                params.append(end)

            order_clause = "ASC" if order.lower() == "asc" else "DESC"
            sql = (
                "SELECT content, confidence, knowledge_type, entities, relationships, insights, source_location, timestamp "
                + "FROM knowledge_items WHERE "
                + " AND ".join(where)
                + f" ORDER BY timestamp {order_clause} LIMIT ?"
            )
            params.append(int(limit))
            cur.execute(sql, tuple(params))
            for row in cur.fetchall():
                (
                    content,
                    confidence,
                    knowledge_type,
                    entities_json,
                    relationships_json,
                    insights_json,
                    source_location,
                    ts,
                ) = row
                results.append(
                    {
                        "content": content,
                        "confidence": float(confidence or 0.0),
                        "knowledge_type": knowledge_type,
                        "entities": json.loads(entities_json or "[]"),
                        "relationships": json.loads(relationships_json or "{}"),
                        "insights": json.loads(insights_json or "[]"),
                        "source_location": source_location,
                        "timestamp": ts,
                    }
                )

        # 메모리 결과 정렬 및 슬라이스
        try:
            results.sort(key=lambda x: x.get("timestamp", ""), reverse=(order.lower() != "asc"))
        except Exception:
            pass

        return {
            "project_id": project_id,
            "query": q,
            "matches": results[: limit],
            "count": len(results[: limit]),
        }
    except Exception as e:
        logger.error(f"지식 검색 실패: {e}")
        raise HTTPException(status_code=500, detail="지식 검색 실패")


@app.post("/api/v1/persuasion")
async def persuasion_from_text(payload: Dict[str, Any] = Body(...)):
    """임의 텍스트로부터 설득력 있는 콘텐츠 생성"""
    try:
        text = str(payload.get("text", ""))
        if not text.strip():
            raise HTTPException(status_code=400, detail="text 필드가 필요합니다")
        source_type = str(payload.get("source_type", "text"))
        # 텍스트를 지식으로 변환 후 설득 컨텐츠 생성
        knowledge = await system._extract_knowledge_from_text(text, source_type)
        persuasive_content = await system._generate_persuasive_content(knowledge)
        return {
            "extracted_knowledge": asdict(knowledge),
            "persuasive_content": asdict(persuasive_content),
            "knowledge_summary": await system._generate_knowledge_summary(knowledge),
            "learning_insights": await system._generate_learning_insights(knowledge),
        }
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"설득 콘텐츠 생성 실패: {e}")
        raise HTTPException(status_code=500, detail="설득 콘텐츠 생성 실패")


if __name__ == "__main__":
    _port = int(
        os.environ.get(
            "ULTIMATE_MEDIA_PORT",
            os.environ.get("MEDIA_KNOWLEDGE_PORT", os.environ.get("PORT", "8001")),
        )
    )
    uvicorn.run(app, host="0.0.0.0", port=_port)
